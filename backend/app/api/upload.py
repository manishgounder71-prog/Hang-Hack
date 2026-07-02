import os
import uuid
import logging
from typing import List

from fastapi import APIRouter, Depends, File, Form, UploadFile, HTTPException, Query
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.database import get_session
from app.core.config import settings
from app.core.cognee_client import remember_content
from app.models.memory import Memory

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/upload", tags=["upload"])

ALLOWED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".csv", ".json", ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".css"}


def _extract_text(file_path: str, filename: str) -> str:
    """Extract text from various file types."""
    ext = os.path.splitext(filename)[1].lower()

    try:
        if ext == ".txt" or ext == ".md":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        elif ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            except ImportError:
                pass

        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(file_path)
                return "\n".join(p.text for p in doc.paragraphs)
            except ImportError:
                pass

        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            except ImportError:
                pass

        elif ext == ".csv":
            try:
                import csv
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    reader = csv.reader(f)
                    return "\n".join([",".join(row) for row in reader])
            except ImportError:
                pass

        elif ext in (".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".css", ".json"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

    except Exception as e:
        logger.warning(f"Error extracting text from {filename}: {e}")

    return f"[File: {filename}] Binary or unsupported format - stored as reference."


@router.post("")
async def upload_file(
    file: UploadFile = File(...),
    user_id: str = Form("default"),
    session: AsyncSession = Depends(get_session),
):
    # Validate extension
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"File type {ext} not supported. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}")

    # Save file
    os.makedirs(settings.UPLOAD_DIR, exist_ok=True)
    file_id = str(uuid.uuid4())
    safe_name = f"{file_id}_{file.filename}"
    file_path = os.path.join(settings.UPLOAD_DIR, safe_name)

    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)

    # Extract text content
    text_content = _extract_text(file_path, file.filename or "unknown.txt")

    # Store in Cognee
    cognee_result = await remember_content(
        content=text_content[:5000],
        content_type=f"file:{ext}",
        user_id=user_id,
        metadata={"filename": file.filename, "file_size": len(content), "file_id": file_id},
        file_path=file_path,
    )

    # Dual persist to SQLAlchemy
    db_memory = Memory(
        user_id=user_id,
        content=text_content[:2000],
        content_type=f"file:{ext}",
        metadata={"filename": file.filename, "file_size": len(content), "file_id": file_id, "file_path": file_path},
        tags=[f"file:{ext}", f"upload"],
    )
    session.add(db_memory)

    return {
        "status": "stored",
        "filename": file.filename,
        "file_size": len(content),
        "content_type": f"file:{ext}",
        "content_preview": text_content[:200],
        "memory_id": db_memory.id,
        "cognee_id": cognee_result.get("cognee_id"),
    }


@router.get("/supported")
async def supported_types():
    return {"allowed_extensions": sorted(ALLOWED_EXTENSIONS)}
